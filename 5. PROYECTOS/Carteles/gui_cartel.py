
from tkinter import Image, Listbox, filedialog, messagebox, ttk, StringVar, TOP
import os
from functools import partial
from base64 import b64decode
import tempfile


from hashlib import sha256

from PIL import Image, ImageTk
import tkinter as tk

import sqlite3

import pptx
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches
from pptx.dml.color import RGBColor
import qrcode
import io

from datetime import datetime
from PIL import Image
from PIL import Image

user_pc = os.environ['USERNAME']
salt = "df^%kTZBnirXv6Yc&#65Ki"
imagen_ventana = ""

def crear_cartel(text_ean, text_ubi, text_la, path_img, text_alternate="MEGAPACK"):

    
    # print("-"*32)
    # print("\nentrada en generacion de carteles")
    # print(f"EAN: {text_ean}")
    # print(f"UBI: {text_ubi}")
    # print(f"LA: {text_la}")
    # print(f"PATH: {path_img}")
    # print(f"ALTERNATE: {text_alternate}")
    
    if text_alternate == "":
        text_alternate = "MEGAPACK"
    
    logo_img = "iVBORw0KGgoAAAANSUhEUgAAAO4AAABBCAYAAADWpJgZAAAAIGNIUk0AAHomAACAhAAA+gAAAIDoAAB1MAAA6mAAADqYAAAXcJy6UTwAAAAGYktHRAD/AP8A/6C9p5MAAAAJcEhZcwAAEnQAABJ0Ad5mH3gAAAAHdElNRQfoCBQNHAlOu87sAAA4RUlEQVR42u1dd5hVxfl+Z067dfsuyy5dFAXFhr1HorGixoIdLDF2kxiNJdEYE2uKPbF3UexRUYModgWxICC9LGXb3d1bT5uZ7/fHuVtAYO8u9t++z3PhuXvPmfOeOfNO+eb7vsPQhz58Bxj3cO3miezqHYgIjG3sSMaiZumCPx36waw9+o34vmn/YKF/3wT68P8DHHwsA64jgHd/NLsnmV39GQD1ffP+oaJPuH34rmACKALAuj2SsZAv3e+b7w8aBfR+fejDNwLKfwo6lnUv7//X6BNuH35wYACjQiX+/xR9wu3DDxJ9ut04+oTbhz78CNEn3D704UeIPuH24TtBTyxTfegefcLtw3eEPtl+k/jR7+MSERYvXYaGxmZYloVhQ4egrKTo+6a1Qa5fLViE1rYUuMbRr6oSQwcN+L5p9eFHiB+dcJ99+ilMfeYR5NIpCCeH4/YYiXfen4c1+d/HlALH77cDuBGCGYkhVFyJOx98FOx72Bi84Oyz4SWbkG5pgHKzOG7PrfH2e3NQn/9930EWxu+7PaBbiERj6Dd0JK77163fdxX34UeAH4VwiQgXnXEKmpfPx5NXnYb4juMsLaSqDSsyVAvFBh42rl8lMc0CEUkpMpKb9bphLVWauSK71RHNJ+09kk48cE9svvW2uPrvd3yrXB99YhLeeelpZBuWw57xJMSoI0LMTFeHTGuI54uBBx+2bwXXtBAjUr4QacGMNVzTl0qur4zvdVrLibOmKRUuxeAttsb1t931fVd9H36g0AHgyEeH9/A0QuC5xrp8J3TvzUZrnfPcSYs3evTtf78Bn773Jo7Yrhab7z8+yqy6HaK1W25jNs/fmZPcXmc0oIRRVJdKp3zBipgSkvnC5y2C2BLvrRs/zFnxL4xwbFblVjvPP+nAz6QWK8dDz7z4jVbkWWf/FtkVn+KJK0/D4P1PLCKjcYxTNHhkpOGL3SKQo3WmakyOqEakQwSOQYox5SvmeoK1CuJLlt336/d9ZnxuWpGZtbuOW3Ly4jmqqGog7njwsR5xOfXpXZF2WqBIobuJhqlZePL4OWv97dfP7YuFTTPgSRtlkRoQGBQpGNxAxksg66VREu6HEVW74p+HvPCN1uO6OOLRzQCwLi1r/e2MiGBoFiqiA3HXuNc2WualU47G0tZ58KSbr58Nt13Kr82f76atAsCfp56J+Y0fYlVqPoTywBlHaaQ/DC0ET7og6nS9ZozB5CEI8tCaWw2hJKJmEUZU7oLNKsbgN3tct9Fr6QCgcQMA4gj8SQsFW1u4BVsfGAAHQGZjB/3urAl45NpLMXLs4fH+W2y7a3rBRxNKZXb/YkuWF1lMj+iAqQEGZ+CMg+UJEIFLgu5JWetI1KY9f682F24q27Dsk2fvfsoIxx4T4bIlpxy6n3/RX67HDtvv0oNbXj9+ddJxaJn1HMzqEUUDdthnT7du9ullZO9ZFFVlRabUo0ZXrqyjGSqACwXdlyqaE2pA1vf3iipy0nbjok8e/9skK176TGTAyEXHHbi3fPK1twvmo3EjpHEjxkkBbOPPReNmDoANBA16dv3b2HnQoXxR8yfFAPS4VR4VJGNSCRY2op4nc8kspaXGdP83+z6cnHffUBpWPhp3HvHWJtfj+qBzkwEoApgBbMSfigg6N5WphVIA5EbvmRlc52aRIqV3Vz/BY8JGy5w4eQxWtM3FVWPvwemTd4swxiIAGAND3CqPmHoknPNShiLV0TswxilqxH2hXKctV58FJBi4Ko/WZP445Tr34PurMLh0JO46cv31qgNAIr0sBMauALAHOrufrje0vu5o3b91VwFdh+nXQXQdAH99B/7quMNQe8BZbFT9yl1YcvVvY4bYq9RS1aVxIGJwGAzgPF8Yw9dGFSKAiEERUBFmcASstCdHJJz0ZU257FFupu2ZUGnlXRcfc1D9RWecgn/d+zB6i/EH7AG+88ksnGzay7SbLihWap/SsKwosQKuOg9M9+082To1pwgAMRQTIBRj/aMIpzy5TcJJbtXSljmxfkbjJDNS9p/jfzamsf+I7fGPu+7pllNLduWRaafpXBX08Bt9LqYWuh/AAyc8MQrDK3fQFyVm7Tjp07/80hW5MYpUeFVyfhyMxUHgjDHHl24rYxAppyn7q0lDZxSHq14rDlV+vP/dUXto2Wjce/QHva7L9SGRWRYD2PUAtsHaU7t17ouYzq0moZzfA9jo8Jh0G6tbsstvcqUzuLv6AdAM0O8ALFn3h9Of3h3LWj5DabhfJGk3b3fQfen9lrZ+voNQXjUAJkliVXJ+XGN6WJEyAcq3AiKAUTPjHpF0pPLTDCBHpMX7y55bvufQ8g9DRvSt/vHhC35298d+RXQAnjpxwVrX1gHAk7YGYGsAe36jtb5h1GEDW1GnHrgbtEhx0dz7Lz4hqnIXVETEVlURhqjBYHBA4wBnDJrGwDVA4wy8S0lEgFKAVAQpCVICBidEDIaSEPRSW41ak8ttnkyuGTN4q9E3lwwZNf2c8YeqOye91KMbOO/K69D2ycsgM1buvHztCcXcuaAqKodXhgOuOu/kxjmDxhFwXoerVICSAVddEUwdiBgMpSHoLY4/ck3Wv7I1642WVuzWnSb85Z3xK5Zg0stvbJSbL50hnsjtoaj7qDil/KmHPlCNolBFvxdm33JW1ktO8KU7tH2K6ElnvecJ5SHjte7viOyEluzqVyNG0WM1RZtNP2FSVjw+/otvrKF4wjYAbAdg127vhctG288UF3DPUU/au3jC3qwACg0AYuv+8YiHB6AoVGmZeuSAL+vf+aUr7AMkiX5EtFa79qULH4VFOkklkfNSezHGjnNEZtnbSye9FLfKnukfH/bBhKdi6sFjZ3UcqwMdI9Z3Gfu43mudEoi2yl2z8KoyzT6lNo5YWYjD0vKNXwNMk8PUGQyDQdcZOA8EzMBAIJAClAoE6wuC7yt4PsH3CTonhDSGIouZ9Vnv4JXZtq0Wv/fSHyv6D3jqhAN29x9//f2Cb2D141ciNvqAat1OXFGm2WcMjLNQqcVh5rnqWsDRNBkMnUHT850NAxhnAAGKAsEqCfhCwfMozzngamkcRSbM1Rnv6NV2cptX/nzsJRQqffm044+V9z/x1IbJUWf9bmyNSwQI5fsV0dqiZS2z/5z12iYGIwO6XRsH5xPzpdvfl+5ET9o//2TV65f1i2/21HlP7+/dfvQb3Rew4XKJ56clLJjKdts285NoudHpdEf1AADJAu9TYp1R+dJXjkSb3RRb0PTh+Tkv+Vuh/IpC6rsH92/40t3cl+5FQnrjlmPOpRWx2heIyG/fHfnBWJVPOnBXGPHyai3TcH2l4YwfEGNWkRWMXLrOYZkMIYvDsjgMs4twtUC8LD8JV5KCEVcQRF64rqvguhT87xHiIJhxDkuTQ5enkzc31PnVVYOG3X7+acPd2+7vftp83L47gsfKao1cw79qQuKw2hiz4gaDrgG63snTNFnQ0Rgcmg5oOgNr50rtHUxeuH6+k3EVnDxX5gdcBxdxhHU5Ymk6eZuflnG/bPsnrjj3VPXXOx7a5HpnjFFzdtW4rNd2oiJl9qThtR+b7wAGpN2WvytSsQvG3PGf2/HGJnpcfJvbd72nRkS4f8Zf+XNzbjsl5TRfpkjGu1p7vgl0qVfmK29Yq91ws6/cJLuY/a/9mB+EcE876mAooJJnGq6sMt0TBsW5ETeCqaVpMETCHFaIIxzOC9fi0E0OzdDANQam8Y59WiKCEgrSl4F4XQXXlfAcBdth0G0F21FgjNA/xqExVb04mftDy6rlifLhOz5yw/U3yEv/cOkGuZ467ucA2ADTbf1b/7AYNyAGI24Eo6llMoTbeYaCj2Fx6AYH1wO+HcIFAymCEhJSKAgv+HiugpXnatsKjhtwrY5yADRoSTp7XXbFTOeY6//5DJUMxN/+eu2mVL0ytXCt7WcOVKRivR0tWDCBgFSiKuelzrv3sz98ePRjwz97+sRF30dz+tZdtKYseGDXjNf2G0Uy/m26B7SvhoXyBjt+7pzjdxw56wnMTQA/AOFefOmVSDfWhfzV835XpjtnDIxxI24CuhaMXJEIRyymIVJkIRSPwojGoGkExgU0TYBrAA8WvgBRIAapoHwF6UkIR8B0BDxbwjAldF1C0xhytgLzFfoFgqhYmMpe07Lks/p/3n7rqxvievKxR8MKGVHVuPg3VYZzwsAY12J5ruEQRzTCEYnpiBaFYMUjMCJRaDoDZy64JsH1PNc8SFHAU8gOrpYrYdoCRk7BMCT0HEMuJwFGqI4yEOTAFRn7mr9f+YemrQ484R1sekPdz5fOoE19ju3WIqm8kSm3+bwBRVtdCCDb6/K0XnvjfqueNsc8trnWklszXki3oD3U9on7+gyo6/v7124mL15P5vZvsxv2OvW5vZ9/6Mi3A+FuctDyJkwVhux7LD6968JDynTv9MFFzCoyA6NOOMQRi2qIF+sorqmAVTEIPFwKbkYBEBg8GKwFjFoA5qPrfpAGAkkF4QoIW0C3feiWD83woWkMGpfgDMjkAJBCZYTBlnLgskz6solHHDB/9cqVS1+bOfdrXAW3uNO4bEKV5pw5MM60qNEp2lhUQ1GZieL+lTDLBwRcjYArJxs6bwNDAmBirafH8xYq4QgIR8DP+dBMD7ouAtFzGTzcbLB86xfhcIUamc2k/jrnjWePn3jsESsfeOr5Xj85T+S2JHwzjmX5RsYckTt4dWbRvwHM7G1ZhhbadELfAjSuFyuSWxO6tx8AHc3SJ1p7aZH/TScqrO6JVEyS2GPs4PEvPoS3VX4fVwcDyzHGUrROD76uJtf5nYiUJZUIEevhupyAiceOx4zbzx8ZkenLq2OqosgMppzto1e8REfZ4IGwqkeARavBrDigBVvNdi6HuUskvpyzGktWrIDjCTAEhqDBFSZ2Gl2CLYdGEYmY8LIedFMH1zm45oHx4MaIgAwRTART0bQndl+dTlxUUjv8Esycu5Yp8NQjD4ElWkdYLHd2bRTxuBmsaUNWINriMgMlQzeDVTUMPFIFmHHACAGkkE6msXgZw9w5dVhWtwquL8AYg6ExDK2ysNM2xRgxNIpwxIBm6dBMDVxzwTW/c6ecAJWVCCGY4meF3GmNm5pYPmK76/H1bbWCumLGwIiCSxTSeRf6gAmqjDG+xZg7MHPmuT1pFJ1FsN4PnN/qVNn2M8MUyYLST+rcXBkxi57KecnPhfJF19844yxilox0/MypQnm13dUtAUyBBr294mkLgK0DQFGoygkbsZs0rj/YXmntd886tEvtBbRrl4ioJpFdeZ6t0jsUXKtBMUxpIGZEuKVy4ypMf3RZONjzNPJr2mhMQ1G/MljVW4DHa4FYP0izGL7S0dKSwfy5LfjyyzSWLnKwaGEKqxJpKAqsyxoHKqaswV47lGLc/tXYelgElqGBaazLWhhQMtg2UlIhqhNqokzPCH+ci/RjE44Z9/GDk7t4BcUqLJlYeFqFJbcszRvNTDOYykdjHEX9KwPRxmuhwpUgMwYfFloTacyZ14g5c9qweIGD+Qta0diWy9dqnmuxgX3HlOGI/fth682iCBk8WLfnrc8AQAqBISunEDUI/aOwWlvtX62e8/Fb5550zDt3PDp5fQ2429bfPhXjTEsBmKdxozVsxDIA4Ipc1BN2hDNtmCIxoNDRAYAJou0v3P3mJ07GxT+1sKAhRFSysQOC0ZWRpUcm/XL0JZfd9s4lIt5lAsEY4AqFo0afoU1d8FBTymm6VpGKdndhRdLMuG0MyK9x026zTDvNn3Q+FepCInj+gXTbnxzDsPLtoita59zsiOwoFNgZ52/I05ixoK3S82Vq+VYxJk6qijA9pHVOO8MRjlixiXBlf/BoNShaDRmpgssiSDRmkGwDoEVRWl6BtrYkKlub0ZrOwfYVGGMgIjS2+XhmWiOmzWzBKQfX4OTDaxEqtjruiSQFxiupQQiCcgglIYYKRw1a4dq/Lh80cjbyHkXnX34tErNeGxlH9rjKMDRLR8fMIBziiBSFYJUPDEbaSCVUtB8cZaG1OYvWFgIzYigtK0N5ZSX6tTQilXPgS1qL61NTGzD9kxacengtTjikBqHifGdJAV8pCb7g8AWBXEKJxVAVkgNW2qkjYgN2eA+93M4LRMub4lbpnwzNelbTwtldB4+TjDHMWvm61pRepses4pEpJ/EHV+QOJQLf2MMOOgJivnK3+LLhwwjy61z6DoxG3wU86ZWjMA9D39CM+StaZosh5YPQlF7RsYdPBJSGS7CwaYasKd78CaG8gxyR3fnr6usEA4POrYzODQLywt229lBsUbkDikMVeWF2dtaUd1Zpl+9rc+9EymkqWZb44gpfuScSUaiwOTrAGPNDevT20ki/O3d5bwJv9l86rsgUW8RNDo0DphUIIRTisIpi4NF+gFUEChXBIQuZtA+u6YgWxeB7PtKpJCKRCEwrDM452tsuY6y9AaE1JXD3syvR0OLholOGIBq3AuOVUAjl90w9j8PzCaYiVIQZq7edg9YsmrPjAePPf/f1SbchWj2Mebr8WaVO1TEjGCVDVmA9Dkc4QvEQeLgEsOJQZhw5YSKXleCagWhxDI7joKi4GNFYDJFIBIbG4Uu5NlcQ6lt93Dm5Dk0tHs46ZiCiRSGQDKzkIb+da7DNZSlCZZihzVcHN87/8v6//OWvc/74xys6qruQRkjBDEVFjKInBxRv8WBzdrUDEBTJjg6DMw316boPaoqG/LE1Vz/ck87IQspWpKrqkguj2AQDVS/xrRqnhLQtgirEcsYlqfDnq9/ELkPG4fJ91w1uaQPwJm5486ImIf2LHZHpB2zswRFMLVRfEat1gbxwbz30WQDPdsvkF/dUoLp4cCTnp8/0pHMOgSKFihaAr3PrqahZflPKbmlsXPphaQTeXqUW0828V5FlMlhWfhslbIEZEUC3YLtAU1saumlCM3RougaucXBNCz6BB8bXn2BeFLav8OybDagqNTHhyBqYURPKVxCuRMhVcF0Ox1WQghAzGUpNWZm1Uzu4pSPeBYA10x6qjErn8OIojHaPKNPkwSekwYiYYLoBaCZyHkdTSwpWJAxN5zBMA4ZpQDeCj6Zp0PjXybZPmzO2xGNT1qC60sIJB/eHETUhPQnTk7BcBcvicF0FIQhxiyHM5eZJO3dILB6fi16MapxpCVOPTp61+n1n5vnB357Fv9Y65vevHIExAw758qGZV7zoS3dLIuLdPXdf2uFVbfOMnvL5oYNABeSFBohId/3sHkWx8imDS7dZuvmNkDELiBiAoYWhayYsLYT3lz+tfOXOVqDZlhaC7afh5wMUunq+MQCmHgXXwwB6kAHj6MdHYJuaffX69JKzXWFf3hPRMoBMzXozZpb8YWnz8vrnJ6xELKQPK7bY5nEz2MnR9bwYLA4jpIEbwU4VKYlVK5qQaGyF63jwHB+u7cLJ2bBzOXiuC993odSG2yxnDJ5PeGzKany5KAsjakIP6zDCOsy8s4RlcnCNweRASQha1GR7jMh+EK6prIJ0sjuFmNguZrAuXINOxgxxaAYH4xxCEFYsa0KyNQ0358BzfNiZHHLZLOxcDk6er1Bq41wF4b7nVuKjL5IwIgaMsAEjpMO0gs7NNDk4BywOFJukadI+YOmXs0p60xA1ri+Lh8oXzDhvw/V308HPY17jh4pz/WPOuFNIuUQEIb2fXHZkznQfYBvtINtv2pfuEU2ZFc8+PPOyf42oqj6nf7xiHICxhmb9LGTE9i0KVe4cNmKjGGObhc3iAUPLR5f2iw4IAWROPSMDV9gYWro1ztn1Npy+2204e/fbccshgWtuQfu4l0w5EqvSK7TP10w9xfbTv1MkS3ow0kLXrFeLQlW/f+HUupUnP7UzQloEws7sbDJZZWnBqGjkxRAIVw+2Q3gWdfWtmDM3g0i8FFII6IYB13HQmkigtaUF6WQb0qkUPLHxJR5jQCIl8OSUNdhyyHBYIR26pUO3fBhG4InVvh0cMxi4Eju0Na0asrqpcR4XA/eIRihu8sBLyzACryjDZNBDBjSDgZGPpcvq8eVCB2UVlRCeB03X4dg2Wpqb0ZJIINnWipa2FFx/41w5Z6hv8fDM1HrsOKoIekSHbuswLA2GqWAYKtjWUoS4yWBmvEFrli0sB9Da45bI2Kq4VZrszup01f73YtzDA1aBsRQIkR5f5yeCsBFtc/y0T4C2sePySzXLk842vnS2YYxLArkApOOnyRVZlXESLkjZiqTLvZSTytW3+NJJKSj35/cWLayKDnqfMT7jmNFnNI29N46pZ6Q7yu9WuBf+9xBUxQZps9e880vbT1+jSPbvgWhJ58Zbpha5POMm5ixs/hybV2yLIVf/zYRwdgybZOp5Meh5X17NYNBMDZqpQKoF73/YigUrBaqq+yGdSkHXdXieh5ZEAo319WisX42GRApSbXyWyFjgy/zOp634ZE4Se29XFFzH4B3uk4EvMSGkM4Q41XDpbfWHy//8VcvHT9fEDDCdB1FJuh6MurrJ82UQCAl8/nkDFq8mZLM5FBUVgWta0Mm0tqJ+zWo0N6xBIpmDKmBCyxnDZ/PTWFJnY8uBVrBFZPDgunrgCy18QlgH4hYv1qLFVQDaXZUKGukYAEuLNI/sv58HTO+2/va+i2cUqf/X7wZh4AsZYy0g1HR7bJd9WyKlAYgwhrwNQUJ22cWTUq61p+f4aXjSbku7ic8OvLfk+bJI7RMH3d/YOOW0JgDdCPc/M28AwNiLX9x8SM5L3qhIdrvf1E4UAHRufF4cqrh4eduaz7atGYPNK7YFAKxe+EVUF+4wKxzM1QMxdBGPrkEzNTQnE/jwozrkRBS5bBpFxSXQdB2+5yHV1opEUwNWrFyDlO3nK2rj5BgDklmBT+alsfeOpXnRaR0BAFr76M8BU6OwBrZFcYRHRDRUEtZdcBasbzt4anlXRpMj4zTjq7mrUZ8Kw85lEIsXQdM0uK6LdLINicZ61K2qR9YtnGtDi4ePZ7dhqyH9wY2gk9F0Bo0H1wcDTI3B4lSck/5gAIVHSeShcS19wS5XqQtxdbfHKlKgHrwe5Ac8T+61hTtqFC9OOY0LJES3wu2oB7bx7+slSIBUokQqsS9jfHeR8beKW2XX/OmVY1dfc/BTGxauL33oXMdhD9Xum/FarxfKH9wTDwudm59GrZILR9Xs+ynhbTx8bKcTTcg044yzCj3vCMEZy0f/sMBJQufQdA1px8f85U1wnSaUNDcjHo9D1w14notUKo1EMouc4+UroxCbQbD90pb2oIC8QwbvvDYPtr04AIODwbH7z5/9RbSCUKbnrQGMBxUfRCWxDq5r6rKYMXsVbJ+jorkJ0VgEXDPguQ6SyTRa0zlkHL/dul7Aw2YQUuHLxRn4CtCMgGvQwXQ+fM4AnZHhe3511+de6HOiIC1DYcf32ENuPWGzP3KUx2rbmnN1b/nK24uItG/LV7lruUTKdEV2IgBamJzz2wdm3WKvV7inTN4Fxm8N/HLM0J2ybstNQvlbFXKx9jwYOjcWWXr4T6/c3vhO7RV1+O+EVWsdl25rjRSTihu8M7C8I9Cc5SNoOIOSCp5QSKR9JHMe9MY2cMYglIKQqou/Z89CWtIpASGp4zodYsxH5jNGgdWOa/GoLiKMZKz9Cgz5jqZ9xOVBR5No9dCSEbA9hVTWhanzIOWLUvB6yxUMLUkPnqdgMpavm3ynwVkX8RKHmy3vTQORyvd7c97/Vzw54zX1sy23elCR3M7xs+OI6FvPQ5h3kjF96R5uu+nJ5z9/4bSvWZWPemQQmtILcfh2NVsnsqv/7klnx/aTuwUBGtfr4qHyS8aNOvuVS649Fvcc9e7XDvNdRwNI6yqGtVy18qkiSCjk3aEgFcH1JWxPwBeqfQ+yV062UhHafUfbI3XawwLbU3QEe6uch3SmcUZrGSLWl3XD94MABwAQipDzAq5uO1fWO65K5MvtmuJrPSMfI9E1RL/g8nXNMFHoQErt67XCyqaf2GgLALce/wCa3JXLyyMDrrD0yBTOeJp6WC+9hSLR31XZYy7Y8yRrLeGe/fzPkXLaEAtVbpb2W2/0lbsnUPicnDOtKaTHrgwbxS+2Oc3qxoPWH+wdLyt3wTRbru9G8xv/RAC39CBKpEvD7/rp3SKKEAlr0DWWz5aRdzBpr3x0xCqAKc9pTHu2IJ7rQg8qfx4RBcH7RCiK6zDNoDrXy7VXIIQjOnRDCzzYqAvPtT3KCUY017trMC3p2wUdqXG917X+U8H4kRPxxmlpPHniV3PLI9UTikKV54eM6HSN62nk+yrq4acQ5EMnme1nd1ncPKN/h3CnLn0Jn636HyJmfFAiU3ezL50DUODMjgjQuNZq6eEbqiIDJrkiK6/Y794NHm9oPKOI2nzVRQhEUIqglIKSCiQVDA5Yeq/DuzZAFogVGeAcIF+B8tkyggD8QISSAEEAMSMRL6uyGedpQZ3nEwXB70pRh2dTVWUI5cVGQRbjQsEADOkfhqkDSigoSZ0+y/lZAwEQBKlbofoupxbMwpNu2fXTzy5oW7A0UhPnXP9hhu2sVWvdHMHWmuP1CjPXfAChRNPLDzQ8NLBk1PjSSP+TdW5eFjaKbg3psZcNzZqlc32xoVkNeT/wLAPr/DBmc8YFQ88ErEjUOCI3VAeAwx+oxs3TTkZlbED/jNt6tSftwwpdeOd9XTMRs+SGLat2vqOhbZn34qkrN36OFU0p8DpfYWeifFyqBITIN1BfQboSlUU6Nh8YRl2DU5jHfLdcCabBMKw6BKYIMn8tJSnI/ZQ30ygCXMV8pVnLhm61TXrFyk8afaszR5SUefG0B+x7EpWlJoYNiGLRKvsbYBpwLYro2GWbYnAEcbtKBN5dUlLH9peQBEE8y3W9bq3zC7kGACHdioa2eRY2kLivHc/PfBAPzr1qIIK3yncLBkac69RNsRs+vRe6Iij4cuOzhwPvL0XKaTUVUU8ymnbghEnbA4zhX++cj0ioEnfcNAHvLJ5cX5+qe8H28cI/DnuRvzrv7siK1JwiX9hxy4gXN2frYp7IGWEjorf7/euaZeiaVZn1UsM9Ye8plbczEfTuNKeUDLfZDRU6EeH4p7ZDEanihszyiz2RO4lQuGgZ4Bha6K6q6JA7GtPLnCdOmtfteTsefrL90X+unGML50hJ4FIF+aGkIAifIH0J4UnEyiwctGclPpydhOOrTU7oowjYrDqMPXYoAwkJ5YngWvk0N0IGyZpcSXAl2sD5vN//7vfitJ+NzOR8giQGpYJjpSTIfLC+9CSsKMPWm8cwbWYC1IMwmg1BKmD4wDC23aoI0pN50ap8ArxglkAEuArI+aotKXIre3MdAm2W9ZI1kz69f8H47U/b4HE10eHMl84oRbKgEVfXLL88OkgCCwo5/BuBUtIQ0gkd9sgg/PfkFV/7/fppv8eypo+wLPPlMKn8XhnzHh//6VrfH10n5Hjsn3+hcAwy/WqQGVAK9CuuglAefCWQzaU6uiMhgbQDLPkD4djHho9uytY97ktv1Mau3e5/L6TH9J/fV4yySE1kVXrZH4Xyfk0goydBA5YWvr84VP63hF2feXkd6/GGsGTJEtJCsY9dL5P2pSo2tXx+KEFBChdXQHoC0vWx905l2PmdYrz5SQt0rfdiICKYOsOxB1ZjYLUFP2VDuBK+p+B7CsIPREEE2D4goC0yrNBCIoGTDwh/lvZzvlAwDNWey4qC81wB4QoYvsCh+1Rh+owWzJqfhqb1mmrAVQP2G1OGoogGP+UFgfbtfP18J0NAzgeyUmtlkWiyN9eSyh+U89rGHrfdxAVTnrwdDx03a63fb3r3YqxqnYP75l5W7fiZsYVaUXVuZGqLh3vfpXDBUMSZPlIn/YPLXz8RfzugM5H8+Me3wvt1j6PYqAg7InskkYr3ZmI07uHBu/nSOTD/dd0UsR1ECMQUyXd/M+HKqT9nF2ywvGt2PQd7Dzlqzovz7nrDl96oQmeWOgOLJLIrzxbKO0OR6knQgDC00CNl0YF/Sdqr23J+Bnvdxdc7KAbGpk4b41uYiNHa4bNtyZbZgraNGAxCEDxPwfM4LFdC2D78kI7i0giOHFuNOYszaE764Lzntd1u7NphRBxjd6sA/CAzhnAFRF4Inh8kmRMKSHlEPvHp5qDtmhh7DRMOCb1vS2elLWiopbN8lI6C73IIN5+9IudjQFUYE8bVYvGdC5HKyk3iuu9OZThibD+QF2TF8B0fvtuZsVIpgq+ApEsQPDRz+KjRTXhhSo+uFXjxqFDWT5127OMj3p2+eMEXxz02AiEtCH/UtBCmfnUfdhgw1lrQNPNYTzoFx11zptVXFw39ziKD8qOR6Qn7AB7ikz+pm5I6/KFa+MKFrpl4a9E87L3ZiJo1mcUTpfKPoF7O33Jecozjp/9USLCBwc3bVn7qT7vy5ZPVtYc8st5jrtrvLvz83jKpSCa6K68zcwaD7viZWiJ1BoGKezDSCoNbr8T04j8nc/PrwcLM0kPQtPVHLShFkErAlwK+AAwTlHS1JuGzuW0ebVsS6hSum0+Wpuf8fDYIF2N3KUNz6yD885GlyDqqR4IgCgS5WW0Yvzt1KPqV6HCTNvycD88WcB3Vkf1RKYItCK0OHFuwWdVOowIAq2rIfL9h3hsp3zm9yAITguC5BNdTMB0JI+cHGTZMDfvsXIaTlvbHfS+sgutTj7kSAVsMDOPXxw5CRZEOpzXg6tsyn60ySOOqCLAFISN4GzMjL+QSq3vtiugJe/vmTN3dYwZU3ipJfOhJagUAg/FwyIgN+HTVG0fn/NRERTLcbaaGfBA559rCQ7f5g3s57sg3tl7uAXQZ0wppn5LEwa25NdfGrLJ7pHLrJAlopFduW1Oye3O27iRP2nsTUY+yWXaFzk0HIEGkuo184kzb4Z3FT9XUZ5eu/NPrEzE/MaMjJQ+RgiIPA4q3wrb99+EPz7i8JB/+3V2ZXjxUkdYVSQ6gRx4gHNwBSLS6jRcGDnB2EG7vtz+jDdc/18CkxCeJC1c+Yd0UfarNEwflfJRojIIUqo6CYyrohoBmeOAaR0jjOGpsFQDCv5+qQ2ObB41vfJuFOrZOCDuOKMKFJw/ByKER+FkXXsaDl/Xg2hJOXrhCEIQCWh1CiqzPY1WDPvrn3Q8AANyVc1yu6y+02HRcRYjFdUZwPAXTYTBNBSMnwA0X3OCwihlO++VAGKaGh15chZa0XzBXANhpZBEumTAUWw4OwUsHXP2cD9cJhOs4gfOJzHO1lT49FI29c+N/HkFvkN/c5560dxHK/U/WS66gfO/PuRZVpPorklU99BJyQPRlm93QK05dETGLletnPEmy22ODe1ExT2TPaVPeIURyFYHIEZlqRXIAEXUERxSarG0911iFQGEbXVIGCd7sMSuT86+tjA66PmIWL/770bP92i72+ztm3sRWNM8IvzDnll085R4EdN+7aUxfU2yWL+rWirW+ypEko1LKo75GtvBinkjUzZhca42bZju5V1od7/iQzhgXBNvJZzc0GDTdB9c4mMZgxi2M/0V/VBSbeGZqPT6Zl0LWkR0xt13R7q/Qv8zEnjuU4oSD+mPLIRG4KQdu2oWb9uDmBBw7EILjBFbarACaXZZReuj+4WNPXoFJQYzyg/+bifGHHjCz1fdntrn+fiGdgfkEx8lH6ugSXPeh6RoYYwgVWTj9qAEYUhvG06+twafz091yrakw8bOdynH8wf0xtL8FL+3CzQvXyUnYOQXbDjoZqYCsT2jIwWnz+CtWsq1X69v2Btb+XBWpmCKvI1Beyk6x9Mg5DSwllLt879oxXdtGr7ZfBhRvmV3ROvsLKXL7FDK9zXdEmlT+MADDAECS6PobGGNS5+ZKobxqIrJ6cm+GFlrOmdasSG3Uup6vTyvnJ09anV64Q2N22Yy3Fj/SxJiWDw0jJpRngWiIr9ztfekOLuA5ka6Zn8StspW9Ss/KNiGTVzsPKw0W1txUTmpPNNg4KG5Sqc4Bx1FBdkNN5l0Q8ydIghlTOHD3Muy6XQnem9WKV99twtylWWSczgbGwVBRrGO30SU4eO9KjBwegyYlnKQDL+PCTblwMz7srEIuJ5GzJXyh4CugKafQbPOZImw8l1j4yVqE9/nva/WfH/mz25rdlu3iJpUWMcB2FDRNBuGAHGDM7VjPG1ETB+5Wjt23LcWbHyfw+ntN+GpZFllv7Te2VZca2HWbEhy6bxW2GBKBpiTcpAMv7cHNuLDTPnIZCTsnYdsKQhI8CTTZoJSypocral588KkXerVio2Ch5lOQbbDb4Phuy8v/b2jWF/3iwxcDn29iMwEYN4TGzWmcORNVgQal7u5DY/qaklDlg4nc6vMIZPWET9wqW560G+Yi3yl0x4OINE/Y23iwt/narIvW9i3rnreWilmlL9x9y7v295ZXmTHgjsmv4szjj34rlVz1YIOdPc/SYTAQcjkVvKqjXbTtuZJFsP0SjRg4eM9y7DOmFPUJD2mXOpOegVAe01FVbkAHwc95cGwfXtYLRq+MgJ2TyOYkstngDQdSAW0OYU2Or/GM+L9EfGDixhvWTjT+/hV3IRqK/y/Rln4xnPVONjXwEAg5u0sHk9+TJqmgfAnpCkQiBsbtU4Gxu5RhdbOLrN+ZQJwRUFmkoarUhAYFP+fCzgVc/YwHJyuRzQhkswrZXLAOFwpocQirsqj39OiNj05+sX59T7xQ9/6QEXtFkRzmCnv0Jj/UwHuuLWwU3bf30PGJx74B4d7/y3cx/vEt32nMuNNdYR+6yQUCZGjWK5wbSxlj4Z7OA8pDtblGbdlHQnm/IKJuZ6xrh/bRBn/vploBAJzr08NabDq9Td9/QnQ/05KxYmX/aky725rc/1ltLMgdlc113rDMJ3ZrTxquZz1oIR2GqWNIudaZEB2dScZFyoHnSQhXwLd9+DkB15awsxK5nEImK/PrRULKI6xIk59WoTurtt71leZVK75Ww4/87RxceNwBmUy49IYmt2VoKOfvXR0JpsyZXPvsJ3DiECJIi2OEg5zOmhkEwQ+rMjoTorMuXNM2XFfCt30IRwRGs6xcq4Nx3IBr0iWsyiKbptCd3Ay/P+nlqRh/yNjePwCiz0Na7HVPun8nUt0anzZcTGDttPTwK2WR/q/Up7t/n+zGiuu62pt061ctvziz7EFfenspkgUZUTfEUeNao6mFJrXk6ocpUj0abQHgtcUvq+2rRz0hMv7+rsjuSwWapzeFMwDo3JgVN0tvOnmbG1sZY9+/cB/87zQAWHHSL3b/S70tqyyNtq6MMMBXQDb/5j2Zd5DwKLDiWho0K4il5UYQoYOuo7Po8nYAV8J3ZGD0chTsnELO7hRC2gfq0iSafXOyUVJ1T3rNUv/Rp55ZL9dbnnwduwDzRh685w2rsmKgyWloeZiBSCGDYL0qBEH4WsDVlTBsrTPuN58ituvo3OF95UsIR8J3A662HXxytoTrUr6DAZanSDR71qNW9fBbE3VLnE0SLcAkCRW1ip+3ReYIX9pjezNlbl9/mVro7fLIwJt8YaevGntvzwrZCO679TpMWzLpNUni346fuUCRCvfUTJ1f26YsPXrryOo93v14xcsjexPZ89Vvga8Gzll82F/6Xy+VP0gqb1ih4u1Rneb/YYDSuflFyIhd9t+J9e9uPyWwu3zvwgWA31x8GWS4fHpq5rOXLs/k7gTU4IowA0hBtjs8eMEWjOVyGKaE0R5UruffHZQvS+VTrwqhIH2C56r8NlNgTLIdBd8PDDxpH1iRUqrB5tNghq/MJVsanujmjX1XpwnPnHvqa8nEst8tTadul6RqKsMM8FSwt+oHjiQhj8N0OExTQjd4kNlDZ2C8kytRMJNQQnU4dLRviQVcgxeBCQWkPcKyFIl6PzJJhov/4idbUs+/9eE3Uf3s5sNnrrno+W1/12rX3yWkt3v7qyILydQPAJxxx9KjU+NW2aVPnjhvbncXLIRT14n+6TtfhsMfqs0Uh/tfz9CQs/30uYpkVYe7wwZ4rs1Ra7H0yG1Rq/zWmw563j/kgcrylJPggZdbIXw6ceeka/Bu3eQ3dM06K2k3XuMJe2eir0eQ9QTrzqIZoBjTWkN69KWoVXzTrkOOntdv0jRcc9D9AALhEgCR/3wXcVgM67zd+583X4czTxxPudItXkPL/IuWp3PXuVJt2S/CESYCkYIvCKanYLkcZnuOKL0zAL+9olSXoAFf5MXgB2+/8/x8fuK848KKNImEpz9NZvjql155d2lzAeQPijP86oSjZWjwtv9tWzrLYtncFb4SW1dFOMIIgiR8n+BYCpaZ52p2ZtjoGkerurzDV/gEzw/2aN28U0g71zaHsDIDp9EPPSqt4qvdVGLV5Lc+2RjNHgX9vj7/bjz3wfIvD9ul5lzHzxzvK+doofzBRLTBvUrGmOTgCV2z3jc08+Wq2JBXHpv12epr3/wVrtzv7q+fEATs+whyNW2snTEAkjG+1jEvnroKhzxQ3TakbPQNSxKffiykd5iCPEiRrKUN+x0TZ7xZ59b/TC305OCykf9bmVxo/3X6BVwoUYLu2zxb3zHn7PEnHPPYFuK5N5ZNPebnm69sza35hSJ1gJDuDgpUBlLG+oIG1g6O/1p9AmCexrSExo2vdG6+xcA+HlI66r2FiU/Ty1s+x0PjO20GOoIX914NoBjfnXAXYR3v83sem4RTjj5C1h7+mxeWv3hrwrftPznS37tfBGaRifzIy+C61PE6y/bcS6xzphxEGkmCyAcC+Hk/ZKkCI5QjgIRDWJNjiSRFH1dFZTewdOOq+x7/N8ad8OuCbuDux5/GqUceLB6a8t6kUw7ao9nN5v7oSrlrVRhmrJ2rZHBdBUPn0PMJ1HW9U7TtUSEyz7XdV7rDD5oAWwAJm1Bvs/pmz3xML625TveyiUc3LtoeP44zd7oEFz57AOY0v//FzkMOnbuwcebjGa9tH1fktmeMD/RELqJIcc40ZemRNBGt0rkx19JDH5WG+305bfFnybiVxue/mYdtqtefc4FAbwE4D93bzRhAX0bN8q8d8/LEepzwRIUrlf/qsJLR05qcVf/2pL2bJ53tOfhgT9pFihQ3NMvjTFtDRHMtPfR2WWTgJ6uS8zJtdhNemdiI0qc+JF+5zxJoXvdckAXwNV/eyScuwC3v/havLnj4q6lTMl+dOn7HB1tya7ZSJEa5wtmCSNRwbpQL5Yd8aZtEnZ5W+cAcZehhT2NaRkivwdTDDYyxRVGjeG7YiC9+5LjPE/vdHULaa8UbZ2a+5gegA0gC2Mhbkr87PPz08zg3FqbUV+++E9nl8AlrMomzUm2ZX/cLU2VZCAgbgCEJggPMY/kUMvn41/YnQJSPl+205EkCPBmMsvU5kglP/9w3S66v3GzUyyvmfpp78n8z8MDrM3rE9aHnXsFF552Pf95269QJh+//VZ2TuaDNdSZUhVVlWZgjrBOMvKGqPbsG2wBXUu2hjcFU3xVAm0dozJFodvXPpBm/loqKX1dO1n7ov1MLYNfzBdctR70OABgweZFIuYnZ7y5unX3e3scaTdlVkWWJzwxHZJmph2l41S7eFiXb56589Saxy9AQLC2MZ06ajp0H7YOnseFEKQxsDoA5hbFhSLuN6/3l8eO/BACUPdTs2X5q9qLmzOwzdz9HW5qYHV3c/InpiRwrDlXLLcrH5G666hlnn19JOH4G03/to91vWoNGDOxdAO8Wxmf9uHDPfwAA/jbtbMyuf7st47Z8sKTV+WDZpYRrpp5utbhNVn16qV6fnM8F+Ws9FF0LU03xlrI6NtSP68X2VQfcJ7e/BSiPeNC4jhvePg9vneUCmIPHj1//rsEPDhed/xukls5EpGZLK7l8zr5R7h8d5f6hxYaqLrKC9KkGB3Tema0CQEffqZA3FCnAEYS0R2hxmUhLfW5O8sddHnnpf6++PeeYk47GnY8+vUlcfz3xVHiNi6FKhoS8+vn7FRnq2KgmDy0xVUXcAGJmwFVrF+8GuPp5rikPaHEgMsr4zFbGJE+PvPjoS28snHDEgXjohdcL4nTIA/1/n3IablCkWHfrVEsL/XnaWc7V6/52w9sX4KvGj+BJH6uT8+AKGyE9isFloxHXixEKl+PGAx/D940/TT0NdW1zIJRAXesc+NJBaaQGNbHhcKSNh4/rWYe8qZi5+kPc89HlMJkOjzw051ajKb0UUvlrbQ2ZehhV8c1QFq6G52VQVjwUZ475K7Yo73Z7GMAPVLjtOOvYw+C2rASr3S7MUyt3JTt1TIS528UNNtRksjysw+gqivagfKEARwKOgO0SX5ny2DJPj73NQvHJm5/29wWz/zWRBm+7O67/+z++Ma5nHHMYZOtKqMoto1qucS/upo4Ow9smZrIhJlNlYR26mffl7srVV4ArAEexrC1ZXcqlZcKMv8UjZc/c//SLi048YDcMGDYCN/z7wYK5fBPC7cMPGz9o4bbjjPFHIrViDkq2O9Tw18wt1Uhs5tuZ7XSmtuEkBzMlSgxdDxEReUJmoJtNCnyhJPaFFQ5/mpN6w2l3vZ584My9qKxfLW558MlvjevEIw9C85LZqBq1u+UnG0ot09xMuLkduXBGalCDGVSxrumWUkq6QmaZbjUKqZYwKzbLikQ+T7S0NUye+n7ypIP3Re1WO+LGf/S8c+mZcMN/nnaWffV38Bj78A3iB7Ed1B3unfQcpr79MZ6+5x8+vFwjJ69R1s/7oN/xN+grPpoSFplWq6K82pLCV81NTW75gFHOMff8O/fwLoNhajUwtDDeueMSPDJlk5Y0BeGB56ZgYWMWN51ztCuFW0/Srm/+/OP3hp5+g9b8xbSo9ByrrLzadF1bNjUnvH4DR9u73X6n/f7YERBUAZMpXP37C/DYlOnAlOmbTqh7/Cg67z6sjR+FcAFg7N47r/X9b3++Cq11s4RFbppBpP1sG5SUsJiE6Sbx9m/Px7m3v4Q9d8578v33ze+M6+ZVa7/q9MarLsPqVV9Ki/yURwJuphVSCoS4hOG2YsG5J2GnY87D+RfmA66nfrSJDH562RX7sDb6etufIA55oPr3Kaexb6r8E8Y3nEKxDz8+9I3OP0b0CfeniT41/sTRJ9w+9OFHiD7h/kTRN+T+tNEn3J8kvu232PTh+8aPZjuoD4WjI4snCpIw+yYSuPfhu0WfcH+C4Iy3csYX09ozqvVJmHHGW79vvn3oOfqE+xNEaajqRR2Ysa5SO1KTdUnZEDXijT0rvQ8/BPwfHR+j4nL8/rEAAAAldEVYdGRhdGU6Y3JlYXRlADIwMjQtMDgtMjBUMTM6Mjg6MDkrMDA6MDBzRknuAAAAJXRFWHRkYXRlOm1vZGlmeQAyMDI0LTA4LTIwVDEzOjI4OjA5KzAwOjAwAhvxUgAAACh0RVh0ZGF0ZTp0aW1lc3RhbXAAMjAyNC0wOC0yMFQxMzoyODowOSswMDowMFUO0I0AAAAASUVORK5CYII="
    img_dcd =b64decode(logo_img)
    img_stream = io.BytesIO(img_dcd)
    user_pc = os.environ['USERNAME']
    
    fecha = datetime.now().date()
    
    #producción
    path_save = f"C:\\Users\\{user_pc}\\GXO\\SPCABANILLAS - Cabanillas_Zooplus\\6. Stock\\10. Plantillas Zooplus\\Carteles LA"
    #pruebas
    #path_save = f"C:\\Users\\{user_pc}\\GXO\\SPCABANILLAS - Cabanillas_Zooplus\\6. Stock\\10. Plantillas Zooplus\\Carteles LA\\Pruebas"
    
    qr = qrcode.QRCode( #type:ignore
        version=1,
        error_correction=qrcode.constants.ERROR_CORRECT_L, #type:ignore
        box_size=10,
        border=4,
    )
    qr.add_data(text_ean)
    qr.make(fit=True)

    qr_img = qr.make_image(fill_color="black", back_color="white").convert('RGB') #type: ignore
    datas = qr_img.getdata()

    #Hacer transparente el fondo de la imagen
    new_data = []

    for item in datas:
        if item[0] == 255 and item[1] == 255 and item[2] == 255:
            new_data.append((255, 255, 255, 0))
        else:
            new_data.append(item)

    qr_img.putdata(new_data)

    # Guardar el QR en un archivo temporal
    usuario = os.environ['USERNAME']
    
    qr_img_path = f"C:\\Users\\{usuario}\\temp_qr_code{fecha}{user_pc}.png"
    qr_img.save(qr_img_path)

    left = Inches(6)
    top = Inches(2)
    height = Inches(3)



    # Crear una presentación de PowerPoint
    prs = Presentation()

    # Añadir una diapositiva en blanco
    slide_layout = prs.slide_layouts[5]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.add_picture(qr_img_path, left, top, height=height)

    # Añadir un rectángulo con bordes redondeados
    left = Inches(1)
    top = Inches(0.5)
    width = Inches(8.3)
    height = Inches(6.3)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)

    # Configurar el borde verde
    line = shape.line
    line.color.rgb = RGBColor(60, 151, 0)  # Color verde para el borde
    line.width = Inches(0.1)  # Ancho del borde

    # Configurar el relleno transparente
    shape.fill.background()

    #Insertamos la imagen del logo
    left = Inches(1.7)
    top = Inches(1)

    slide.shapes.add_picture(img_stream, left, top)

    #Insertamos el texto de la presentación (UBICACION)
    left = Inches(5.35)
    top = Inches(0.75)
    width = Inches(2)
    height = Inches(1.7)

    textbox = slide.shapes.add_textbox(left, top, width, height)

    text_frame = textbox.text_frame
    text_location = text_frame.add_paragraph()
    text_location.text = text_ubi
    text_location.font.size = Inches(0.65)
    text_location.font.color.rgb = RGBColor(0, 0, 0)
    text_location.font.bold = True

    #Insertamos el texto de la presentación (LA)
    left = Inches(5)
    top = Inches(4.75)
    width = Inches(2)
    height = Inches(2)

    textbox = slide.shapes.add_textbox(left, top, width, height)

    text_frame = textbox.text_frame
    text_location = text_frame.add_paragraph()
    text_location.text = "LA " + text_la
    text_location.font.size = Inches(0.75)
    text_location.font.color.rgb = RGBColor(0, 0, 0)
    text_location.font.bold = True
    text_location.font.underline = True

    #Insertamos el texto de la presentación (Texto altenativo)
    left = Inches(1.55)
    top = Inches(4.75)
    width = Inches(2)
    height = Inches(2)

    textbox = slide.shapes.add_textbox(left, top, width, height)

    text_frame = textbox.text_frame
    text_location = text_frame.add_paragraph()
    text_location.text = text_alternate
    text_location.font.size = Inches(0.75)
    text_location.font.color.rgb = RGBColor(151, 12, 0)
    text_location.font.bold = True
    #text_location.font.underline = True

    #añadir la foto
    left = Inches(2.45)
    top = Inches(2.55)
    width = Inches(2)
    height = Inches(2)


    img_path = path_img
    slide.shapes.add_picture(img_path, left, top, height=height)


    save_data = path_save + f"\\{fecha}-{text_la}.pptx"
    # Guardar la presentación
    prs.save(save_data)
    os.remove(qr_img_path)
    
#Copiar al portapapeles de la ventana de busqueda
def copiar_al_portapapeles(event):
    widget = event.widget
    app2.clipboard_clear()
    app2.clipboard_append(widget.get("1.0", "end-1c"))

# Funcion para actualizar las etiquetas de la primera ventana y de la segunda
def actualizar_etiqueta(id_articulo, descrip_articulo, mensaje, ean_articulo):
    
    label4.config(text=id_articulo)
    label6.config(text=descrip_articulo)
    label8.config(text=mensaje)
    label_ean_main.config(text=ean_articulo)
    
    if app2:
        label_desc_text.config(text=descrip_articulo)
                
        # Habilitar el widget Text, insertar el EAN y deshabilitarlo nuevamente
        label_ean_text.configure(state="normal")
        label_ean_text.delete("1.0", "end")  # Borrar contenido previo
        label_ean_text.insert("1.0", ean_articulo)
        label_ean_text.tag_configure("center", justify="center")
        label_ean_text.tag_add("center", "1.0", "end")
        label_ean_text.configure(state="disabled")

def actualizar_etiqueta_ubicacion(ubicacion):
    #print(f"actualizar: {ubicacion}")
    label12.config(text=ubicacion)
    
# Función para abrir el cuadro de diálogo y seleccionar archivos de imagen
def seleccionar_archivo():
    global archivo
    app2.attributes("-topmost", False)
    
    archivo = filedialog.askopenfilename(
        title="Selecciona una imagen",
        filetypes=[("Archivos de imagen", "*.jpg;*.jpeg;*.png;*.gif;*.bmp")]
    )
    if archivo:
        label10.config(text="Cargada")
        cargar_y_mostrar_imagen(archivo)
        app2.attributes("-topmost", True)
        return archivo

def cargar_y_mostrar_imagen(ruta):
    #print(ruta)
    # Cargar la imagen usando PIL
    imagen = Image.open(ruta)
    
    # Redimensionar la imagen a 200x200 píxeles (puedes ajustar este tamaño)
    imagen = imagen.resize((200, 200), Image.LANCZOS)
    
    # Convertir la imagen a un formato que Tkinter pueda usar
    imagen_tk = ImageTk.PhotoImage(imagen)
    
    # Mostrar la imagen en la segunda ventana
    etiqueta_imagen.config(image=imagen_tk)
    etiqueta_imagen.image = imagen_tk  # type: ignore Guardar una referencia para evitar que Python elimine la imagen

# recupera la imagen en binario desde la bd y la transforma en un tmp en jpg
def recuperar_imagen_desde_db(db_path, la):
    date = datetime.now().date()
    conexion = sqlite3.connect(db_path)
    cursor = conexion.cursor()

    cursor.execute("SELECT imagen FROM imagenes WHERE la = ?", (la, ))
    imagen = cursor.fetchone()
    name_tmp = f'{date}_{la}_tmp.jpg'
    #C:\Users\jgmeras\GXO\SPCABANILLAS - Cabanillas_Zooplus\6. Stock\08. Ideas, Sugerencias, Cambios\BBDD --NO TOCAR--
    path= os.path.join('Users', user_pc, 'GXO', 'SPCABANILLAS - Cabanillas_Zooplus', '6. Stock', '08. Ideas, Sugerencias, Cambios','BBDD --NO TOCAR--', name_tmp)    
    
    with open(path, 'wb') as archivo:
        archivo.write(imagen[0])  # Escribir la imagen en el archivo
    
    conexion.close()

    return imagen

#****IMPLEMENTAR LA BUSQUEDA DE IMAGENES   ****
def busqueda(la_search):
     
    
    if not la_search:
        messagebox.showerror("Error", "No se ha introducido ningún artículo")
        return None
    
    #db_path = f"C:\\Users\\{user_pc}\\GXO\\SPCABANILLAS - Cabanillas_Zooplus\\6. Stock\\08. Ideas, Sugerencias, Cambios\\BBDD --NO TOCAR--\\img_database.db"
    #print(db_path)
    conexion = sqlite3.connect(f"C:\\Users\\{user_pc}\\GXO\\SPCABANILLAS - Cabanillas_Zooplus\\6. Stock\\08. Ideas, Sugerencias, Cambios\\BBDD --NO TOCAR--\\artículos.db")
    cursor = conexion.cursor()
    
    
    cursor.execute("SELECT * FROM articulos WHERE la = ?", (la_search,))
    
    # ***** REVISAR DESDE AQUÍ ******
    
    #imagen = recuperar_imagen_desde_db(db_path, la_search)    #aqui nos quedamos
    
    
    resultado = cursor.fetchone()
    
    # Verificar si se obtuvo algún resultado
    if resultado:
        # Asignar los valores a variables
        id_articulo, la_articulo, descrip_articulo, ean_articulo = resultado
        conexion.close()
        return id_articulo, la_articulo, descrip_articulo, ean_articulo#, imagen
        
    else:
        conexion.close()
        return None
    
def seleccionar_ubicacion():
    app3 = tk.Toplevel(app)
    app3.title("Seleccionar Ubicación")
    app3.geometry("500x300")  # Ajustado para más espacio
    
    def limitar_caracteres(entry_text, limite):
        limite = int(limite)
        if len(entry_text) > limite:
            messagebox.showerror("Error", f"El texto no puede superar los {limite} caracteres")
            app3.focus_force()
            return False
        
        return True
    
    
    
    def habilitar_input():
        if check_var.get() == 1:
            entry.config(state='normal')
            menu_desplegable.config(state='disabled')
            menu_desplegable_pasillo.config(state='disabled')
            menu_desplegable_columna.config(state='disabled')
            menu_desplegable_nivel.config(state='disabled')
            menu_desplegable_ubicacion.config(state='disabled')
        elif check_var.get() == 2:
            entry.config(state='disabled')
            menu_desplegable.config(state='normal')
            menu_desplegable_pasillo.config(state='normal')
            menu_desplegable_columna.config(state='normal')
            menu_desplegable_nivel.config(state='normal')
            menu_desplegable_ubicacion.config(state='normal')
        else:
            entry.config(state='disabled')
            menu_desplegable.config(state='disabled')
            menu_desplegable_pasillo.config(state='disabled')
            menu_desplegable_columna.config(state='disabled')
            menu_desplegable_nivel.config(state='disabled')
            menu_desplegable_ubicacion.config(state='disabled')
    
    def buscar_ubicacion_bbdd(menu_desplegable):
        seccion_seleccionada = menu_var.get()
        
        if not seccion_seleccionada:
            messagebox.showerror("Error", "No se ha introducido ninguna sección")
            return None 
        
        conexion = sqlite3.connect(f"C:\\Users\\{user_pc}\\GXO\\SPCABANILLAS - Cabanillas_Zooplus\\6. Stock\\08. Ideas, Sugerencias, Cambios\\BBDD --NO TOCAR--\\localizaciones.db")
        cursor = conexion.cursor()

        cursor.execute("SELECT DISTINCT pasillo FROM localizaciones WHERE seccion = ?", (seccion_seleccionada,))
        pasillo = [str(p[0]) for p in cursor.fetchall()]
        
        cursor.execute("SELECT DISTINCT columna FROM localizaciones WHERE seccion = ?", (seccion_seleccionada,))
        columna = [str(c[0]).zfill(3) for c in cursor.fetchall()]
        
        cursor.execute("SELECT DISTINCT nivel FROM localizaciones WHERE seccion = ?", (seccion_seleccionada,))
        nivel = [str(n[0]).zfill(2) for n in cursor.fetchall()]
        
        cursor.execute("SELECT DISTINCT ubicacion FROM localizaciones WHERE seccion = ?", (seccion_seleccionada,))
        ubicacion = [str(u[0]) for u in cursor.fetchall()]
        
        conexion.close()

        if pasillo:
            return pasillo, columna, nivel, ubicacion
        else:
            return None
    
    # Checkbutton para seleccionar habilitar entrada o lista
    check_var = tk.IntVar(value=0)
    check_button1 = ttk.Radiobutton(app3, text="Habilitar entrada manual", variable=check_var, value=1, command=habilitar_input)
    check_button1.place(x=20, y=40)
    
    check_button2 = ttk.Radiobutton(app3, text="Habilitar desplegable", variable=check_var, value=2, command=habilitar_input)
    check_button2.place(x=200, y=40)
    
    #limitar caracteres
    limite_caracteres = 12
    validacion = app3.register(limitar_caracteres)
    entry_text = tk.StringVar()
    
    # Entry que inicialmente estará deshabilitado
    entry = ttk.Entry(app3, state='disabled', textvariable=entry_text, validate='key')
    entry['validatecommand']=(validacion, '%P', limite_caracteres)
    entry.place(x=60, y=80)
    
    
    
    # Menú desplegable principal
    opciones = ['Elija Opción', 'BIN', 'SUELO', 'RIAB', 'RACK']
    menu_var = tk.StringVar()
    menu_var.set(opciones[0])
    menu_desplegable = ttk.OptionMenu(app3, menu_var, *opciones)
    menu_desplegable.config(state='disabled')
    menu_desplegable.place(x=60, y=120)

    # Buscar ubicaciones al seleccionar una opción en el primer menú desplegable
    def actualizar_submenus(*args):
        ubicaciones = buscar_ubicacion_bbdd(menu_desplegable)
        if ubicaciones:
            pasillo, columna, nivel, ubicacion = ubicaciones
            
            # Actualizar menús desplegables
            menu_var_pasillo.set(pasillo[0] if pasillo else "")
            menu_desplegable_pasillo['menu'].delete(0, 'end')
            for option in pasillo:
                menu_desplegable_pasillo['menu'].add_command(label=option, command=tk._setit(menu_var_pasillo, option))

            menu_var_columna.set(columna[0] if columna else "")
            menu_desplegable_columna['menu'].delete(0, 'end')
            for option in columna:
                menu_desplegable_columna['menu'].add_command(label=option, command=tk._setit(menu_var_columna, option))

            menu_var_nivel.set(nivel[0] if nivel else "")
            menu_desplegable_nivel['menu'].delete(0, 'end')
            for option in nivel:
                menu_desplegable_nivel['menu'].add_command(label=option, command=tk._setit(menu_var_nivel, option))

            menu_var_ubicacion.set(ubicacion[0] if ubicacion else "")
            menu_desplegable_ubicacion['menu'].delete(0, 'end')
            for option in ubicacion:
                menu_desplegable_ubicacion['menu'].add_command(label=option, command=tk._setit(menu_var_ubicacion, option))
    
    menu_var.trace('w', actualizar_submenus)
    
    def seleccionar():
        
        if check_var.get() == 2:
            pass_pasillo = menu_var_pasillo.get()
            pass_columna = menu_var_columna.get()
            pass_nivel = menu_var_nivel.get()
            pass_ubicacion = menu_var_ubicacion.get()
            pass_local = pass_pasillo + ". " +pass_columna + pass_nivel + pass_ubicacion
            #print(f"seleccionar: {pass_local}")
            actualizar_etiqueta_ubicacion(pass_local)
            app3.destroy()
        elif check_var.get() == 1:
            
            pass_local = entry_text.get()
            actualizar_etiqueta_ubicacion(pass_local)
            app3.destroy()
    
    # Menús desplegables para pasillo, columna, nivel y ubicación
    
    tk.Label(
        app3,
        text="Pasillo:",
        font=("Arial", 12)
    ).place(x=60, y=160)
    
    menu_var_pasillo = tk.StringVar()
    menu_desplegable_pasillo = ttk.OptionMenu(app3, menu_var_pasillo, '')
    menu_desplegable_pasillo.config(state='disabled')
    menu_desplegable_pasillo.place(x=60, y=200)
    
    tk.Label(
        app3,
        text=".",
        font=("Arial", 18, "bold")
    ).place(x=140, y=190)
    
    tk.Label(
        app3,
        text="Columna:",
        font=("Arial", 12)
    ).place(x=160, y=160)
    
    menu_var_columna = tk.StringVar()
    menu_desplegable_columna = ttk.OptionMenu(app3, menu_var_columna, '')
    menu_desplegable_columna.config(state='disabled')
    menu_desplegable_columna.place(x=160, y=200)
    
    tk.Label(
        app3,
        text="Nivel:",
        font=("Arial", 12)
    ).place(x=260, y=160)
    
    menu_var_nivel = tk.StringVar()
    menu_desplegable_nivel = ttk.OptionMenu(app3, menu_var_nivel, '')
    menu_desplegable_nivel.config(state='disabled')
    menu_desplegable_nivel.place(x=260, y=200)
    
    tk.Label(
        app3,
        text="Ubicación:",
        font=("Arial", 12)
    ).place(x=360, y=160)
    
    menu_var_ubicacion = tk.StringVar()
    menu_desplegable_ubicacion = ttk.OptionMenu(app3, menu_var_ubicacion, '')
    menu_desplegable_ubicacion.config(state='disabled')
    menu_desplegable_ubicacion.place(x=360, y=200)
    
    button_salir = tk.Button(
        app3,
        text="Salir",
        foreground="red",
        command=app3.destroy,
        width=18,
    )
    button_salir.place(x=300, y=250)
    
    button_seleccionar = tk.Button(
        app3,
        text="Seleccionar",
        command=seleccionar,
        width=18,
    )
    button_seleccionar.place(x=100, y=250)

#Revisar cuando esté ubicaciones listo
def generar_cartel():
    #print("*"*32)
    #print("entramos en la funcion generar cartel")
    
    text_la = label4.cget("text")
    #print(f"LA: {text_la}")
    text_alternate = label8.cget("text")
    #print(f"Mensaje: {text_alternate}")
    text_ubi = label12.cget("text")
    #print(f"Ubicacion: {text_ubi}")
    ean_safe = label_ean_main.cget("text")
    #print(f"EAN: {ean_safe}")
    path_img = archivo
    #print(f"PATH: {path_img}")
    
    #print("asignadas variables")
    #print("*"*32+"\n")
    
    try:
        if text_alternate == None:
            #print("sin texto")
            crear_cartel(ean_safe, text_ubi, text_la, path_img)
            messagebox.showinfo("Cartel Generado", "El cartel se ha generado con éxito")
        else:
            #print(f"texto: {text_alternate}")
            crear_cartel(ean_safe, text_ubi, text_la, path_img, text_alternate)
            messagebox.showinfo("Cartel Generado", "El cartel se ha generado con éxito")
    except Exception as e:
        print("generar_cartel")
        messagebox.showerror("Error", f"No se ha podido generar el cartel \n error {e}")

def seleccionar_archivo_creacion():
    global archivo2
    ventana_crear_la.attributes("-topmost", False)
    
    archivo2 = filedialog.askopenfilename(
        title="Selecciona una imagen",
        filetypes=[("Archivos de imagen", "*.jpg;*.jpeg;*.png;*.gif;*.bmp")]
    )
    if archivo2:
        ventana_crear_la.attributes("-topmost", True)
        
        print(archivo2) # test
        
        return archivo2


def crear_LA():
        
    def crear_la3(entry_la, entry_desc, entry_EAN):
        ruta_imagen = archivo2
        conexion = sqlite3.connect(f"C:\\Users\\{user_pc}\\GXO\\SPCABANILLAS - Cabanillas_Zooplus\\6. Stock\\08. Ideas, Sugerencias, Cambios\\BBDD --NO TOCAR--\\artículos.db")
        cursor = conexion.cursor()

        conexion2 = sqlite3.connect(f"C:\\Users\\{user_pc}\\GXO\\SPCABANILLAS - Cabanillas_Zooplus\\6. Stock\\08. Ideas, Sugerencias, Cambios\\BBDD --NO TOCAR--\\img_database.db")
        cursor2 = conexion2.cursor()



        la = entry_la.get()
        desc = entry_desc.get()
        ean = entry_EAN.get()

        cursor.execute('''
                        INSERT INTO articulos (la, descripcion, ean)
                        VALUES (?, ?, ?)''', (la, desc, ean))
        conexion.commit()
        conexion.close()

        cursor2.execute('''
                        INSERT INTO imagenes (la, imagen)
                        VALUES (?, ?)''', (la, ruta_imagen))
        conexion2.commit()
        conexion2.close()

        messagebox.showinfo("Éxito", "LA creada correctamente")
    global ventana_crear_la
    ventana_crear_la = tk.Toplevel(app)
    ventana_crear_la.title("Crear LA")
    ventana_crear_la.geometry("400x300")
    tk.Label(
        ventana_crear_la,
        text="Nueva LA",
    ).place(x=40, y=30)
    entry_la = ttk.Entry(ventana_crear_la)
    entry_la.place(x=200, y=30)
    entry_la.focus()
        
    tk.Label(
        ventana_crear_la,
        text="Descripción",
    ).place(x=40, y=70)
    entry_desc = ttk.Entry(ventana_crear_la)
    entry_desc.place(x=200, y=70)
        
    tk.Label(
        ventana_crear_la,
        text="EAN",
    ).place(x=40, y=110)
    entry_EAN = ttk.Entry(ventana_crear_la)
    entry_EAN.place(x=200, y=110)

    tk.Label(
        ventana_crear_la,
        text="Seleccionar Imagen",
    ).place(x=40, y=150)
    entry_img = ttk.Entry(ventana_crear_la)
    entry_img.place(x=200, y=150)
    
    seleccionar_img = tk.Button(
        ventana_crear_la,
        text="Seleccionar Imagen",
        command=seleccionar_archivo_creacion)
    seleccionar_img.place(x=200, y=180)
    button_crear_la = tk.Button(
        ventana_crear_la,
        text="Crear",
        command=partial(crear_la3, entry_la, entry_desc, entry_EAN),
        width=10,
    )
    button_crear_la.place(x=100, y=250)
    
    button_salir = tk.Button(
        ventana_crear_la,
        text="Salir",
        foreground="red",
        command=ventana_crear_la.destroy,
        width=10,
    )
    button_salir.place(x=300, y=250)
    

def crear_usuario():
    
    def crear_usuario3(entry_usuario, entry_password, entry_rep_password):
        if entry_password.get() == entry_rep_password.get():
            conexion = sqlite3.connect(f"C:\\Users\\{user_pc}\\GXO\\SPCABANILLAS - Cabanillas_Zooplus\\6. Stock\\08. Ideas, Sugerencias, Cambios\\BBDD --NO TOCAR--\\usuarios.db")
            cursor = conexion.cursor()
            
            usuario = entry_usuario.get() + salt
            password = entry_password.get() + salt
            
            user_data = usuario.encode('utf-8')
            password_data = password.encode('utf-8')
            
            user_hash = sha256(user_data).hexdigest()
            pass_hash = sha256(password_data).hexdigest()
            
            cursor.execute('''
                        INSERT INTO usuarios (usuario, contraseña)
                        VALUES (?, ?)''', (user_hash, pass_hash))
            conexion.commit()
            conexion.close()
            
            messagebox.showinfo("Éxito", "Usuario creado correctamente")
        else:
            messagebox.showerror("Error", "Las contraseñas no coinciden")
        
    
    ventana_crear_usuario = tk.Toplevel(app)
    ventana_crear_usuario.title("Crear Usuario")
    ventana_crear_usuario.geometry("400x200")
    
    tk.Label(
        ventana_crear_usuario,
        text="Nuevo Usuario",
    ).place(x=40, y=30)
        
    tk.Label(
        ventana_crear_usuario,
        text="Nueva Contraseña",
    ).place(x=40, y=70)
        
    tk.Label(
        ventana_crear_usuario,
        text="Repetir Contraseña",
    ).place(x=40, y=110)
    
    button_salir = tk.Button(
        ventana_crear_usuario,
        text="Salir",
        foreground="red",
        command=ventana_crear_usuario.destroy,
        width=10,
    )
    button_salir.place(x=300, y=150)
    
    entry_usuario = ttk.Entry(ventana_crear_usuario)
    entry_usuario.place(x=200, y=30)
    entry_usuario.focus()
    
    entry_password = ttk.Entry(ventana_crear_usuario, show="*")
    entry_password.place(x=200, y=70)
    
    entry_rep_password = ttk.Entry(ventana_crear_usuario, show="*")
    entry_rep_password.place(x=200, y=110)
    
    button_crear = tk.Button(
        ventana_crear_usuario,
        text="Crear",
        foreground="green",
        command=partial(crear_usuario3, entry_usuario, entry_password, entry_rep_password),  # Pasamos los campos correctos
        width=18,
    )
    button_crear.place(x=150, y=150)
    
def modificar_la():
    
    def busqueda_datos(la_busqueda):
        la_search = la_busqueda.get()
        
        id, la, descripcion, ean = busqueda(la_search) #type: ignore
        if la:
            label_descrip.config(text=descripcion, background="white", width=30)
            label_ean.config(text=ean, background="white", width=30)
            
    
    
    def habilitar_input():
        if check_var.get() == 1:
            entry_descrip.config(state='normal')
            entry_ean.config(state='disabled')
        elif check_var.get() == 2:
            entry_ean.config(state='normal')
            entry_descrip.config(state='disabled')
        elif check_var.get() == 3:
            entry_ean.config(state='normal')
            entry_descrip.config(state='normal')
        else:
            entry_ean.config(state='disabled')
            entry_descrip.config(state='disabled')

    def modificar(check_var, entry_descrip, entry_ean):
        if check_var.get() == 1:
            descripcion = entry_descrip.get()
            conexion = sqlite3.connect(f"C:\\Users\\{user_pc}\\GXO\\SPCABANILLAS - Cabanillas_Zooplus\\6. Stock\\08. Ideas, Sugerencias, Cambios\\BBDD --NO TOCAR--\\artículos.db")
            cursor = conexion.cursor()
            cursor.execute("UPDATE articulos SET descripcion = ? WHERE la = ?", (descripcion, entry_LA.get()))
            conexion.commit()
            conexion.close()
            messagebox.showinfo("Éxito", "Descripción modificada correctamente")
            app_modificar_la.focus()
        elif check_var.get() == 2:
            ean = entry_ean.get()
            conexion = sqlite3.connect(f"C:\\Users\\{user_pc}\\GXO\\SPCABANILLAS - Cabanillas_Zooplus\\6. Stock\\08. Ideas, Sugerencias, Cambios\\BBDD --NO TOCAR--\\artículos.db")
            cursor = conexion.cursor()
            cursor.execute("UPDATE articulos SET ean = ? WHERE la = ?", (ean, entry_LA.get()))
            conexion.commit()
            conexion.close()
            messagebox.showinfo("Éxito", "EAN modificado correctamente")
            app_modificar_la.focus()
        elif check_var.get() == 3:
            descripcion = entry_descrip.get()
            ean = entry_ean.get()
            conexion = sqlite3.connect(f"C:\\Users\\{user_pc}\\GXO\\SPCABANILLAS - Cabanillas_Zooplus\\6. Stock\\08. Ideas, Sugerencias, Cambios\\BBDD --NO TOCAR--\\artículos.db")
            cursor = conexion.cursor()
            cursor.execute("UPDATE articulos SET descripcion = ?, ean = ? WHERE la = ?", (descripcion, ean, entry_LA.get()))
            conexion.commit()
            conexion.close()
            messagebox.showinfo("Éxito", "Descripción y EAN modificados correctamente")
            app_modificar_la.focus()
    
    app_modificar_la = tk.Toplevel(app)
    app_modificar_la.title("Modificar LA")
    app_modificar_la.geometry("600x400")


    check_var = tk.IntVar(value=0)
    check_button_mod_la1 = ttk.Radiobutton(app_modificar_la, text="Descripcion", variable=check_var, value=1, command=habilitar_input)
    check_button_mod_la1.place(x=500, y=40)
    
    check_button_mod_la2 = ttk.Radiobutton(app_modificar_la, text="EAN",variable=check_var, value=2, command=habilitar_input)
    check_button_mod_la2.place(x=500, y=80)
 
    check_button_mod_la3 = ttk.Radiobutton(app_modificar_la, text="Ambos",variable=check_var, value=3, command=habilitar_input)
    check_button_mod_la3.place(x=500, y=120)
 
    
    tk.Label(
        app_modificar_la,
        text="Zona de modificación de datos",
        font=("Consolas", 18, "bold"),
    ).place(x=70, y=30)
    
    
    tk.Label(
        app_modificar_la,
        text="LA a modificar",
    ).place(x=40, y=70)
    
    tk.Label(
        app_modificar_la,
        text="Descripción",
    ).place(x=40, y=110)
    
    label_descrip = tk.Label(
        app_modificar_la,
        text="",
    )
    label_descrip.place(x=200, y=110)
    
    tk.Label(
        app_modificar_la,
        text="Nueva Descripción",
    ).place(x=40, y=150)
    
    tk.Label(
        app_modificar_la,
        text="EAN",
    ).place(x=40, y=190)
    
    label_ean = tk.Label(
        app_modificar_la,
        text="",
    )
    label_ean.place(x=200, y=190)
    
    tk.Label(
        app_modificar_la,
        text="Nuevo EAN",
    ).place(x=40, y=230)
    
    entry_LA = tk.Entry(
        app_modificar_la,
        width=30,
        font=("Consolas", 12),
        bg="white",
        
    )
    entry_LA.place(x=200, y=70)
    
    entry_descrip = tk.Entry(
        app_modificar_la,
        width=30,
        font=("Consolas", 12),
        bg="white",
        state="disabled",
    )
    entry_descrip.place(x=200, y=150)
    
    entry_ean = tk.Entry(
        app_modificar_la,
        width=30,
        font=("Consolas", 12),
        bg="white",
        state="disabled",
    )
    entry_ean.place(x=200, y=230)
    
    button_modificar = tk.Button(
        app_modificar_la,
        text="Modificar",
        foreground="green",
        command=partial(modificar, check_var, entry_descrip, entry_ean),
        width=18,
    )
    button_modificar.place(x=300, y=350)
    
    button_salir = tk.Button(
        app_modificar_la,
        text="Salir",
        foreground="red",
        command=app_modificar_la.destroy,
        width=18,
    )
    button_salir.place(x=450, y=350)
    
    button_buscar = tk.Button(
        app_modificar_la,
        text="Buscar",
        command=partial(busqueda_datos, entry_LA),
        width=18,
    )
    button_buscar.place(x=100, y=350)
    
    # *** crear zona de muestreo de imagenes ***
    
def crear_usuario2():
    
    def habilitar_campos():
        user_salt = usuario_mod.get() + salt
        user_data = user_salt.encode('utf-8')
        user_hash = sha256(user_data).hexdigest()

        conexion = sqlite3.connect(f"C:\\Users\\{user_pc}\\GXO\\SPCABANILLAS - Cabanillas_Zooplus\\6. Stock\\08. Ideas, Sugerencias, Cambios\\BBDD --NO TOCAR--\\usuarios.db")
        cursor = conexion.cursor()
           
        cursor.execute('''
                   SELECT usuario
                   FROM usuarios
                   WHERE usuario = ?''', (user_hash,))
        resultado = cursor.fetchone()
        
        conexion.close()    
        
        if resultado:
            nueva_contraseña.config(state="normal")
            re_nueva_contraseña.config(state="normal")
            return user_hash
        else:
            messagebox.showerror("Error", "Usuario no encontrado")
            return None
    
    def actualizar_contraseña():
        user_hash = habilitar_campos()
        if user_hash is not None:
            if nueva_contraseña.get() == re_nueva_contraseña.get():
                conexion = sqlite3.connect(f"C:\\Users\\{user_pc}\\GXO\\SPCABANILLAS - Cabanillas_Zooplus\\6. Stock\\08. Ideas, Sugerencias, Cambios\\BBDD --NO TOCAR--\\usuarios.db")
                cursor = conexion.cursor()

                contraseña = nueva_contraseña.get() + salt
                pass_data = contraseña.encode('utf-8')
                pass_hash = sha256(pass_data).hexdigest()

                cursor.execute('''
                               UPDATE usuarios
                               SET contraseña = ?
                               WHERE usuario = ?''', (pass_hash, user_hash))
                conexion.commit()
                conexion.close()

                messagebox.showinfo("Éxito", "Contraseña modificada correctamente")
            else:
                messagebox.showerror("Error", "Las contraseñas no coinciden")

    # Configuración de la ventana para crear/modificar usuario
    app_crear_usuario = tk.Toplevel(app)
    app_crear_usuario.title("Modificar Contraseña")
    app_crear_usuario.geometry("450x200")
    
    tk.Label(
        app_crear_usuario,
        text="Usuario",
    ).place(x=40, y=30)
    
    tk.Label(
        app_crear_usuario,
        text="Nueva Contraseña",
    ).place(x=40, y=70)
    
    tk.Label(
        app_crear_usuario,
        text="Repetir Contraseña",
    ).place(x=40, y=110)
    
    usuario_mod = tk.Entry(app_crear_usuario)
    usuario_mod.place(x=170, y=30)
    usuario_mod.focus()
    
    nueva_contraseña = tk.Entry(
        app_crear_usuario,
        show="*",
        state="disabled"
    )
    nueva_contraseña.place(x=170, y=70)
    
    re_nueva_contraseña = tk.Entry(
        app_crear_usuario,
        show="*",
        state="disabled"
    )
    re_nueva_contraseña.place(x=170, y=110)
    
    button_comprobar_usuario = tk.Button(
        app_crear_usuario,
        text="Comprobar Usuario",
        command=habilitar_campos
    )
    button_comprobar_usuario.place(x=300, y=30)
    
    button_actualizar_contraseña = tk.Button(
        app_crear_usuario,
        text="Actualizar Contraseña",
        command=actualizar_contraseña,
        state="disabled"
    )
    button_actualizar_contraseña.place(x=150, y=150)

    def on_user_verified():
        # Si el usuario fue verificado, habilita el botón para actualizar la contraseña
        if habilitar_campos():
            button_actualizar_contraseña.config(state="normal")
    
    button_comprobar_usuario.config(command=on_user_verified)
    
def modificar_articulos():
    mostrar_ventana_login()
    
# Función para iniciar sesión
def iniciar_sesion(username, password):
    conexion = sqlite3.connect(f"C:\\Users\\{user_pc}\\GXO\\SPCABANILLAS - Cabanillas_Zooplus\\6. Stock\\08. Ideas, Sugerencias, Cambios\\BBDD --NO TOCAR--\\usuarios.db")
    cursor = conexion.cursor()
        
    user_salt = username + salt
    pass_salt = password + salt
    
    user_data = user_salt.encode('utf-8')
    pass_data = pass_salt.encode('utf-8')
    
    user_hash = sha256(user_data).hexdigest()
    pass_hash = sha256(pass_data).hexdigest()
    
    cursor.execute('''
                   SELECT *
                   FROM usuarios
                   WHERE usuario = ? AND contraseña = ?''', (user_hash, pass_hash))
    
    resultado = cursor.fetchone()
    conexion.close()
    
    return resultado is not None

# Función para modificar datos en la base de datos
def modificar_datos_bbds():
    app_datos = tk.Toplevel(app)
    app_datos.title("Modificacion de Bases de Datos")
    app_datos.geometry("500x300")
    
    tk.Label(
        app_datos,
        text="Opciones disponibles", 
        font=("Arial", 14, "bold")
    ).place(x=60, y=30)
    
    button_crear_usuario = tk.Button(
        app_datos,
        width=18,
        text="Crear Usuario",
        command=crear_usuario,
    )
    button_crear_usuario.place(x=40, y=110)
    
    button_modificar_datos = tk.Button(
        app_datos,
        width=18,
        text="Modificar Datos de LA",
        command=modificar_la,
    )
    button_modificar_datos.place(x=40, y=70)
    
    button_modificar_usuario = tk.Button(
        app_datos,
        width=18,
        text="Modificar contraseña",
        command=crear_usuario2,
    )
    button_modificar_usuario.place(x=40, y=150)
    
    button_crear_LA = tk.Button(
        app_datos,
        width=18,
        text="Crear LA",
        command=crear_LA,
    )
    button_crear_LA.place(x=210, y=70)
    
    button_salir_datos = tk.Button(
        app_datos,
        width=18,
        text="Salir",
        foreground="red",
        command=app_datos.destroy
    )
    button_salir_datos.place(x=110, y=250)

# Función para validar los datos de inicio de sesión
def validar_datos(event=None):
    username = entry_usuario.get()
    password = entry_passwd.get()
    
    if iniciar_sesion(username, password):
        ventana_login.destroy()  # Cerrar la ventana de validación de credenciales
        modificar_datos_bbds()
    else:
        messagebox.showerror("Error", "Usuario o contraseña incorrectos")
        ventana_login.lift()
# Función para mostrar la ventana de inicio de sesión
def mostrar_ventana_login():
    global ventana_login, entry_usuario, entry_passwd
    
    ventana_login = tk.Toplevel(app)
    ventana_login.title("Acceso a Zona Restringida")
    ventana_login.geometry("400x200")
    
    tk.Label(
        ventana_login,
        text="Usuario:"
    ).place(x=40, y=30)
    entry_usuario = tk.Entry(ventana_login)
    entry_usuario.place(x=100, y=30)
    entry_usuario.focus()
    
    tk.Label(
        ventana_login,
        text="Password:"
    ).place(x=40, y=70)
    
    entry_passwd = tk.Entry(ventana_login, show="*")
    entry_passwd.place(x=100, y=70)
    
    boton_validar = tk.Button(
        ventana_login,
        text="Entrar",
        command=validar_datos,  # Aquí llamas a validar_datos
        width=18
        )
    boton_validar.place(x=80, y=110)
    
    button_salir = tk.Button(
        ventana_login,
        text="Salir",
        command=ventana_login.destroy,
        width=18,
        foreground="red"        
        
    )
    button_salir.place(x=80, y=150)
    
    entry_usuario.bind("<Return>", validar_datos)
    entry_passwd.bind("<Return>", validar_datos)
#VENTANAS
  
def seleccionar_articulo():
    global ean_safe
      
    def buscar_datos(event=None):
        #global imagen_ventana
        la_search = str(text_entrada.get())
        mensaje = mensaje_text.get("1.0", "end-1c")
        try:
            id, la, descripcion, ean = busqueda(la_search) #type: ignore
        
            if la:
                actualizar_etiqueta(la, descripcion, mensaje, ean)
                #imagen_ventana = imagen
            else:
                actualizar_etiqueta("", "", "Artículo no encontrado", ean)
                #messagebox.showwarning("Aviso", "Artículo no encontrado, revisar numeración de LA")
        except:
            messagebox.showerror("Error", "No se ha encontrado ningún artículo")
    global app2
    app2 = tk.Toplevel(app)
    app2.attributes("-topmost", True)
    app2.title("Seleccionar Artículo")
    app2.geometry("400x720")
    
    icon_zooplus = ""

    icon_path2 = os.path.join(tempfile.gettempdir(), "app2_icon.ico")
    with open(icon_path2, "wb") as icon2_file:
        icon2_file.write(b64decode(icon_zooplus))
    
    app2.iconbitmap(icon_path2)
        
    current_dir_var = tk.StringVar()

    global etiqueta_imagen
    etiqueta_imagen = ttk.Label(app2)
    etiqueta_imagen.place(x=80, y=420)
    
    label = ttk.Label(
        app2,
        text="Seleccionar Artículo",
        foreground="black",
        font=("Arial", 18, "bold"),
        
    )
    label.place(x=85, y=10)
    
    global text_entrada
    text_entrada = ttk.Entry(
        app2,
        textvariable=current_dir_var,
        width=50
    )
    text_entrada.place(x=50, y=100)
    text_entrada.bind("<Return>", buscar_datos)
    text_entrada.focus()
   

    button_indicadores = tk.Button(
        app2,
        text="Buscar",
        command=buscar_datos,
        width=18
    )
    button_indicadores.place(x=50, y=50)

    button_salir = tk.Button(
        app2,
        text="Salir",
        command=app2.destroy,
        width=18,
        foreground="red"        
        
    )
    button_salir.place(x=230, y=670)
    
    label_descripcion = ttk.Label(
        app2,
        text="Descripción:",
        font=("Helvetica", 12, "bold")
        )
    label_descripcion.place(x=50, y=150)

    global label_desc_text
    label_desc_text = tk.Label(
        app2,
        text="",
        background="white",
        width=40,
        height=2
)
    label_desc_text.place(x=50, y=180)
    
    label_ean = ttk.Label(
        app2,
        text="EAN:",
        font=("Helvetica", 12, "bold")
        )
    label_ean.place(x=50, y=220)

    global label_ean_text
    label_ean_text = tk.Text(
        app2,
        bd=0,
        background="white",
        width=40,
        height=2,
        font=("TkDefaultFont", 10)
)
    label_ean_text.place(x=50, y=250)
    label_ean_text.configure(state="disabled")
    label_ean_text.bind("<Control-c>", copiar_al_portapapeles)
        
    
    label_mensaje = ttk.Label(
        app2,
        text="Mensaje (Opcional):",
        font=("Helvetica", 12, "bold")
    )
    label_mensaje.place(x=50, y=290)
    
    global mensaje_text
    mensaje_text = tk.Text(
        app2,
        width=40,
        height=2,
        font=("TkDefaultFont", 10)
    )
    mensaje_text.place(x=50, y=320)
    
    busqueda_img = ""
    
    if imagen_ventana:
        image = tk.PhotoImage(data=b64decode(busqueda_img))
    else:
        image = imagen_ventana
    label2 = ttk.Label(
        app2,
        image=image,
        
    )
    label2.place(x=240, y=300)
    
    
    button_sel_photo = tk.Button(
        app2,
        text="Seleccionar Imagen",
        command=seleccionar_archivo,
        width=18
    )
    button_sel_photo.place(x=200, y=50)
    app2.mainloop()
    
    
app = tk.Tk()

#imagen icono ventana
icon_main ="AAABAAEAICAAAAEAIACoEAAAFgAAACgAAAAgAAAAQAAAAAEAIAAAAAAAgBAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQAAAA0AAAAlAAAAHgAAAAkAAAABAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAABAAAABYAAANACChPkBtZlNYaToLEBR8/gAAAAD4AAAAZAAAABgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACAAAACMCDhxaDjpqrh5pq+4piM//NJjb/0+v6f8/mNf/KHKw7RJBc7gEGS9wAAAANQAAABQAAAAEAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACAAAAEAAAADMGHTt5LWGOyjaFxPotkNb/LZLY/y2S2P81mdz/VLPr/06u6P9Jqeb/Q6Ti/zWMzf5FgK7kEzxnqgIPH2AAAAAsAAAADwAAAAMAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAEAAAAGgADCkgNMVuaRoKu4XrD6P6X4/7/VrLn/zGV2v8vk9n/LZLY/zaa3P9Yt+3/U7Lq/02t6P9IqeX/Rqfj/5fj//+R3fr/aa/a/BlZldkLLVicAAkPUgAAACIAAAAEAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQMbN0ATQ3W2LHi480ij3v+X4///jtz8/5Xh/v9Xsuj/NJjb/zKW2v8wlNn/OJvd/1267/9Xtu3/UrHq/0ys5/9Kqub/itn6/5Pg/f+V4f7/NJjb/y6P1P8hdLr5FVGNxgckUCMAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACJ2ysvEyr5v9Kqub/UK/o/5De/P98z/X/asHv/z2f3/84m93/NZnc/zOX2/87nt//Yb7y/1y57/9Wtez/UbDq/0ys5/9Zter/bsXx/3fL8/84m93/M5fb/y+U2f8ohsz/Cz13QgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAItdrS/UrHq/0+u6P9RsOn/XLjs/0ip5f9CpOL/P6Hh/zqa2f8zisf/NJLS/z+h4f9lwfT/UaHS/1Cj1v9Tr+f/UK/p/0qr5v9FpuT/QKLh/zue3/82mtz/Mpba/yqHzf8LPXdCAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAjB4tb9Xtu3/VLPr/1Gw6f9Nrej/Sqvm/0Sh3P87kMn/O4vC/zOCu/81jMj/Q6Tj/2jC8/9Omcj/QZTL/0CSyf9Jm9D/S6fg/0mq5v9EpeP/P6Hh/zqd3v81mdz/LInO/ws9d0IAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACMnm2v1y67/9Zt+7/VrXs/06o4P9FmM7/TqLW/2S97v9tx/b/QZHH/zWGv/9HqOT/ZLro/0uazP8uktf/LZLY/zKT1v86k9D/QZfP/0Sf2v9DpOP/PqDg/zmc3v8vi9D/Cz13QgAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAQze7fAYb7x/1iw4/9NoNP/T6bb/1257v9kwfP/asb2/2/K+P9sxPP/PYzC/0ml3/9gsN7/NpTV/y2S2P8vk9n/Mpfb/zeb3f89n+D/QqDd/0Od1/9Andn/PZ/f/zGN0f8JMmFRAAAABwAAAAEAAAAAAAAAAAAAAAQAAAAWAAADQixqo9lOo9f/SaPc/06u6P9WtOz/Xbrv/2PA8/9pxfX/bsn4/3LN+v9nvu3/SqHX/0if1v8tktj/LpPY/zKW2v82mtz/PJ7f/0Kj4v9IqeX/Tq7o/1W07P9UsOf/RqHd/xdIfsQEGDFxAAAAKwAAAAMAAAAEAw8hRA45aq0ha63vNJLV/z+h4f9Gp+T/Ta3o/1Sz6/9bue//Yr/y/2jE9f9uyfj/csz6/3TO+/9bs+X/LpLX/y6S2P8xldr/Npnc/zue3/9Bo+L/R6jl/02t6P9Us+v/Wrju/2G+8f9nwvT/a8b2/1ep3/4xc6/ZAC9fEAAkWw4cZ7DxLZDW/zOX2/84m93/PqDg/0Wm5P9MrOf/U7Lr/1q47v9hvvL/Z8P1/23I9/9YtOb/LovE/w1sqv8Oa6f/F3Gq/yuLyv86nd7/QKLh/0an5P9MrOf/U7Lr/1q37v9gvfH/ZsL0/2vG9v9vyvj/csz5/zF1tbAAAAAAAAAAABRSmmMph83+Mpba/zea3f89oOD/RKXj/0ur5v9Sser/Wbfu/1257v9Cn9b/HXy4/w1trf8Nb6//Dm+v/wtopf8DVYX/BFWF/w5jmP8lgrz/QaDb/1Kx6v9Yt+3/X7zw/2XB8/9qxfb/b8n4/3LM+v9LmdHtCTp/GgAAAAAAAAACAAAADhZVmZYvj9T/Nprc/zyf3/9CpOL/Sqrm/0im4P8ti8T/EG+r/wxsq/8Nb6//DnGz/w9ztf8loHH/LrVo/xuKaf8LW2P/BEBf/wZLd/8JWo//G3ax/zqa1f9bue3/acX2/27J+P9yzPr/Zbvs/xVGf3QIIEofAAAAAAw8eRURTJO9EkmGlhldpOc0l9r/O57f/zKRzv8ceLL/Cmah/wtppv8MbKv/Dm+w/w9ytf8Qdbn/EHa7/yuwc/8/4m7/Pd9s/zvaaf8xv1v/D1Y3/wc+U/8JSHH/D2+v/xN7w/8rktb/Ta/p/23I9/8sdbf7FFKZ3hNOloYAAAAAAAAAAAs5fxYUU5plGF2luiJzuvoYcq7/CGCX/wljnf8KZ6P/DGup/w1vr/8PcrT/EHW5/xF4vv8SesH/Kq1z/zzda/8722n/Oddm/ympTv8RXSz/BzlO/wlLd/8PcrX/EnvC/xWBy/8bfcf/G2y10xRVnW4PPIcRAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADUOGExJPmmMZZKq3E2an+hBtq/8Mbav/DnCy/w90uP8ReL3/EnvC/xN9xv8oqXX/Othn/znWZv830mP/G4E6/wxLOf8FM1H/CU57/xN0tv8Xbrb1GWGqoBFJjjsAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAOOH8SElGYYRZhprUTaK75EnW6/xF5wP8TfMX/FH/J/yamdf830mP/NtBh/zTNX/8qr1D/EFwt/wg1T/8LNFjmEEN/fA88eBEAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA8/fxASTpdeFGGpshNttvcVfsj/JKF1/zTLXv8zyV3/Mcdb/zDDWf8agDn/DUwn3QAAAH4AAABNAAAABwAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAADN3DxJQmV8jl1/hMMRZ/zDDWf8vwFf/Lb1V/ymzT/8UbC//Cz8bwAAAAHkAAAA8AAAABwAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAACSiSK0bfjj/KK9O+iy6U/8quFH/KbVP/ySnSP8UbS//DUsg0gERB4UAAABQAAAAKxJnK28PaiowAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAFm4xLhFcKCwZgjlQJq1L+iiyTf8nsEv/Ja1J/yOnRv8ahzn/FnAw/RRoK+QZfzTkGYQ4ngAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAZfTVHIaBE6iSrSP8jqUf/IqdF/yKlRP8ipUT/IJ5B/BmFN4oAfwACAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAANXSgTGoc5fB2UPskdlj7eHZM91xyMOZoQcTAvAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAAA//////////////////gf///gB///gAD//AAAH/AAAAfgAAAH4AAAB+AAAAfgAAAH4AAAB+AAAAGAAAAAAAAAAAAAAAGAAAABAAAAAQAAAAGAAAAD8AAAH/4AAH//wAB///gAP//+AAf//gAP///AD///4D////////////////8="

#imagen en grande ventana
raw_image = "iVBORw0KGgoAAAANSUhEUgAAAHwAAAAuCAMAAADdho1wAAAAV1BMVEVHcEz/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/OgD/maxbAAAAHHRSTlMAOZ4DJi3nBvftGwy7zNVo3RNcRk20w4eoe5JvWvIB4QAAA/5JREFUWMPNWNuypCoMbS94Q1QUUdT//85x79Y2IRG7ps6pmlQ/AWZ1FrmR1+tfk0RqrWX+H2nL5dfKdNQuo1VK2TF2pkqunUpckiJleQq2JNa2vrXVsRuqB5OjSZXbJUU3959Phq78SNbDr/rs2rEV0lYAZU08yHvsKIbIhzTuUJfHYFUBM9IGIPQhbUVtbvjXEwP9I6OkINt0GQj/VHwq1y7jdBVxyppteeitWM8jPSCxjM5VA1abU7UYb7RtzUCxTXd3ev5cFCL+5KNSDOlRs91K2RPs7O5sJ25ut32vTQzpUbdt36Pfnz4xKPHdL8VRSUlPmy0oJWK+stsz6YT4ZU8CcqSky3F7kAZ4XbL45pZZVhDSf06mgKLCvF4rQ7rzte3iLcVXxA0oxgrlBpGKqI87GLaHtECP1aKj5ogMh7YzUWRcjfB//vfB5YwcrNVvI/fbcBNJSRIedjHj6YhGaw4N0lgudyCf2SziOaExCY2FdJ5MIsNjkAgr+LeLgfmrSjyWnrYI+hC88VmjFArR42OtuQsCXhDxhHRpeaf2Q7B5czJA7pIvqr1gksLHfQW4w9X/dCW8r1zCDsoaCNweuC4p4BUw3f2uAI+10LmHtb1kNYATkkVARE7kWqEs3m5ec3VyjwGc7ev8fg/mDGBKyzirp1Aq9rhvHQJ/rcVNtsxHGk3wX1/fKek5O8hnfo7E4FrBvQUk/4vHgnEg4I6N9sDNLbEeuEEJuRFfgic+OEu7rrcQeKW80pdztJtn2lmHc1sInFbB677mQJijQHwrBP455lyZI+CGJNguZUJtTh5DDVqZvS8vd/X4lowD90lHwQZsy9JQklmJHe64uVNiBpyQjoiH6XUKhPkRiFXHd4v4Si5wUwaqWgILSxYFCoumTf+on8BRq/zQRNnqthwunP/E+gEckF442rPDUN6/AtdeQSM/JUyiBqeOLicF4X6CQ9LnvKevFUTk1hxNmd9GXZmhR6FTxqaS+f6qFqsioQZJ36+UtNIkORZq6gfSQILanXu9SdHZcaybksZ5At8nzu/oDH3DnL2ztwLjQDRPbf4BDkm32m8sjtaI5mZfsCOa8itwSPrR7iGkg3jRhXV1OAST/gn9p7tDpB9AKBWfxA9ZSFVGCn0fPF/8NuBDyTUQEzOwGAK2d0yTMajA+T73npNXHUtQqj/5uB017KHM9aTVckN9tqTE0+HjFQ0sziJe8UOWctJ3s6iZ+aBbBL3cDNYAFKmfSRGnrZwDvXkeTdYbhbXp8WLT4/5uPgX3CaIBWw5qQ6MwNUUPs0Ap+mmurVJ2XtYBDAFlCsR7vFZgC84NPyNFOy5tpL94De2/XO6SvP5CEm6Y+rfK/lf5A5RpN2/mmB9oAAAAAElFTkSuQmCC"
icon_path = os.path.join(tempfile.gettempdir(), "app_icon.ico")
with open(icon_path, "wb") as icon_file:
    icon_file.write(b64decode(icon_main))

app.geometry("550x400")
app.title("Generador de Carteles")

app.iconbitmap(icon_path)

label = ttk.Label(
    app,
    text="Generador de Carteles",
    foreground="black",
    font=("Helvetica", 20, "bold"),
)
label.place(x=60, y=40)

button_sel_arti = tk.Button(
    app,
    text="Seleccionar artículo",
    command=seleccionar_articulo,
    width=18
   
)
button_sel_arti.place(x=50, y=100)

button_sel_ubi = tk.Button(
    app,
    text="Seleccionar Ubicación",
    command=seleccionar_ubicacion,
    width=18
   
)
button_sel_ubi.place(x=50, y=140)

button_generar_cartel = tk.Button(
    app,
    text="Generar Cartel",
    command=generar_cartel,
    width=18
    
)
button_generar_cartel.place(x=50, y=180)

button_operaciones = tk.Button(
    app,
    text="Modificar BBDD",
    command=modificar_articulos,
    width=18
   
)
button_operaciones.place(x=50, y=220)

raw_image_data = b64decode(raw_image)
with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_image_file:
    temp_image_file.write(raw_image_data)
    temp_image_path = temp_image_file.name

image = Image.open(temp_image_path)
photo = ImageTk.PhotoImage(image)

label2 = tk.Label(app, image=photo, text="")#type: ignore
label2.place(x=50, y=310)

button_salir = tk.Button(
    app,
    text="Salir",
    command=app.destroy,
    width=12,
    foreground="red"
    
)
button_salir.place(x=425, y=350)

label3 = tk.Label(
    app,
    text="LA: ",
    foreground="black",
    font=("Helvetica", 13, "bold"),
)
label3.place(x=260, y=100)

label4 = tk.Label(
    app,
    text="",
    background="white",
    width=30,
    height=2
)
label4.place(x=300, y=100)

label5 = tk.Label(
    app,
    text="Descrip: ",
    foreground="black",
    font=("Helvetica", 13, "bold"),
)
label5.place(x=220, y=140)

label6 = tk.Label(
    app,
    text="",
    background="white",
    width=30,
    height=2
    
    
)
label6.place(x=300, y=140)

label7 = tk.Label(
    app,
    text="Mensaje: ",
    foreground="black",
    font=("Helvetica", 13, "bold"),
)
label7.place(x=220, y=180)

label8 = tk.Label(
    app,
    text="",
    background="white",
    width=30,
    height=2
    
    
)
label8.place(x=300, y=180)

label9 = tk.Label(
    app,
    text="Foto: ",
    foreground="black",
    font=("Helvetica", 13, "bold"),
)
label9.place(x=240, y=220)

label10 = tk.Label(
    app,
    text="",
    background="white",
    width=30,
    height=2
    
    
)
label10.place(x=300, y=220)

label11 = tk.Label(
    app,
    text="Ubicación: ",
    foreground="black",
    font=("Helvetica", 13, "bold"),
)
label11.place(x=200, y=260)

label12 = tk.Label(
    app,
    text="",
    background="white",
    width=30,
    height=2
    
    
)
label12.place(x=300, y=260)

label_ean_main = tk.Label(
    app,
    text="",
    
)
label_ean_main.place(x=600, y=1000)

app.mainloop()